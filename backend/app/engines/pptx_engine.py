from pathlib import Path
from PIL import Image

from app.engines.pdf_engine import EngineError


def images_to_pptx(image_paths: list[Path], output_path: Path):
    from pptx import Presentation
    from pptx.util import Emu

    if not image_paths:
        raise EngineError("No images to convert.", "QD-PPT-NONE")
    if len(image_paths) > 200:
        raise EngineError("Maximum 200 slides per presentation.", "QD-PPT-LIMIT")

    prs = Presentation()
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for idx, path in enumerate(image_paths):
        try:
            with Image.open(path) as im:
                img_w, img_h = im.size
        except Exception:
            raise EngineError(f"Could not read image: {path.name}", "QD-PPT-READ")
        if img_w <= 0 or img_h <= 0:
            raise EngineError(f"Invalid image dimensions: {path.name}", "QD-PPT-READ")

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        fit_scale = min(slide_w / img_w, slide_h / img_h)
        draw_w = int(img_w * fit_scale)
        draw_h = int(img_h * fit_scale)
        left = int((slide_w - draw_w) / 2)
        top = int((slide_h - draw_h) / 2)

        slide.shapes.add_picture(str(path), left, top, width=draw_w, height=draw_h)

    prs.save(output_path)
    if output_path.stat().st_size == 0:
        raise EngineError("Failed to build the presentation.", "QD-PPT-OUT")
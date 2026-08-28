from typing import Optional, List, Tuple
import io
from pathlib import Path

import fitz
from PIL import Image

from app.engines.pdf_engine import EngineError


def _add_text_watermark_to_page(page: fitz.Page, text: str, position: str, opacity: float, size: str, rotation: int, style: str = "single"):
    rect = page.rect
    font_sizes = {"small": 24, "medium": 48, "large": 72}
    font_size = font_sizes.get(size, 48)
    if style == "tile":
        shape = page.new_shape()
        step_x = max(font_size * 6, rect.width / 4)
        step_y = max(font_size * 4, rect.height / 4)
        for y in range(0, int(rect.height) + int(step_y), int(step_y)):
            for x in range(0, int(rect.width) + int(step_x), int(step_x)):
                morph = None
                if rotation:
                    morph = (fitz.Point(x, y), fitz.Matrix(rotation))
                shape.insert_text(
                    fitz.Point(x, y),
                    text,
                    fontsize=font_size,
                    color=(0, 0, 0),
                    stroke_opacity=opacity,
                    fill_opacity=opacity,
                    morph=morph,
                )
        shape.commit()
        return
    positions = {
        "top-left": (rect.x0 + 36, rect.y0 + 36),
        "top-center": (rect.x0 + rect.width / 2, rect.y0 + 36),
        "top-right": (rect.x1 - 36, rect.y0 + 36),
        "center": (rect.x0 + rect.width / 2, rect.y0 + rect.height / 2),
        "bottom-left": (rect.x0 + 36, rect.y1 - 36),
        "bottom-center": (rect.x0 + rect.width / 2, rect.y1 - 36),
        "bottom-right": (rect.x1 - 36, rect.y1 - 36),
    }
    x, y = positions.get(position, positions["center"])
    shape = page.new_shape()
    morph = None
    if rotation:
        morph = (fitz.Point(x, y), fitz.Matrix(rotation))
    shape.insert_text(
        fitz.Point(x, y),
        text,
        fontsize=font_size,
        color=(0, 0, 0),
        stroke_opacity=opacity,
        fill_opacity=opacity,
        morph=morph,
    )
    shape.commit()


def _add_image_watermark_to_page(page: fitz.Page, img_bytes: bytes, position: str, opacity: float, size: str, rotation: int, style: str = "single"):
    rect = page.rect
    with Image.open(io.BytesIO(img_bytes)) as im:
        w, h = im.size
    size_factors = {"small": 0.15, "medium": 0.25, "large": 0.4}
    factor = size_factors.get(size, 0.25)
    img_w = rect.width * factor
    img_h = img_w * h / w
    if style == "tile":
        step_x = img_w + 24
        step_y = img_h + 24
        try:
            pix = fitz.Pixmap(img_bytes)
            for y in range(0, int(rect.height), int(step_y)):
                for x in range(0, int(rect.width), int(step_x)):
                    page.insert_image(fitz.Rect(x, y, x + img_w, y + img_h), pixmap=pix, overlay=True, rotate=rotation)
        except Exception:
            for y in range(0, int(rect.height), int(step_y)):
                for x in range(0, int(rect.width), int(step_x)):
                    page.insert_image(fitz.Rect(x, y, x + img_w, y + img_h), stream=img_bytes, overlay=True, rotate=rotation)
        return
    positions = {
        "top-left": (rect.x0 + 36, rect.y0 + 36),
        "top-center": (rect.x0 + rect.width / 2 - img_w / 2, rect.y0 + 36),
        "top-right": (rect.x1 - 36 - img_w, rect.y0 + 36),
        "center": (rect.x0 + rect.width / 2 - img_w / 2, rect.y0 + rect.height / 2 - img_h / 2),
        "bottom-left": (rect.x0 + 36, rect.y1 - 36 - img_h),
        "bottom-center": (rect.x0 + rect.width / 2 - img_w / 2, rect.y1 - 36 - img_h),
        "bottom-right": (rect.x1 - 36 - img_w, rect.y1 - 36 - img_h),
    }
    x, y = positions.get(position, positions["center"])
    try:
        pix = fitz.Pixmap(img_bytes)
        page.insert_image(fitz.Rect(x, y, x + img_w, y + img_h), pixmap=pix, overlay=True, rotate=rotation)
    except Exception:
        page.insert_image(fitz.Rect(x, y, x + img_w, y + img_h), stream=img_bytes, overlay=True, rotate=rotation)


def watermark_pdf(input_path: Path, output_path: Path, watermark_type: str,
                  text: str = "", image_bytes: Optional[bytes] = None,
                  position: str = "center", opacity: float = 0.3,
                  size: str = "medium", rotation: int = 0,
                  apply_to: str = "all", selected_pages: Optional[List[int]] = None,
                  style: str = "single"):
    doc = fitz.open(input_path)
    try:
        n = doc.page_count
        target_pages = list(range(n)) if apply_to == "all" else [p - 1 for p in (selected_pages or []) if 1 <= p <= n]
        if apply_to == "selected" and not target_pages:
            raise EngineError(
                "No valid pages selected. Check your page numbers and try again.",
                "QD-WATERMARK-PAGES",
            )
        for i in target_pages:
            page = doc[i]
            if watermark_type == "text":
                _add_text_watermark_to_page(page, text, position, opacity, size, rotation, style)
            else:
                _add_image_watermark_to_page(page, image_bytes, position, opacity, size, rotation, style)
        doc.save(output_path, garbage=4, deflate=True)
    finally:
        doc.close()


def watermark_presentation(input_path: Path, output_path: Path, watermark_type: str,
                           text: str = "", image_bytes: Optional[bytes] = None,
                           position: str = "center", opacity: float = 0.3,
                           size: str = "medium", rotation: int = 0,
                           apply_to: str = "all", selected_slides: Optional[List[int]] = None,
                           style: str = "single"):
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(input_path)
    n = len(prs.slides)
    target_indices = list(range(n)) if apply_to == "all" else [i - 1 for i in (selected_slides or []) if 1 <= i <= n]
    if apply_to == "selected" and not target_indices:
        raise EngineError(
            "No valid slides selected. Check your slide numbers and try again.",
            "QD-WATERMARK-SLIDES",
        )

    size_factors = {"small": 0.15, "medium": 0.25, "large": 0.4}
    factor = size_factors.get(size, 0.25)
    sizes = {"small": 18, "medium": 36, "large": 54}
    font_size = sizes.get(size, 36)
    pos_map = {
        "top-left": (0, 0),
        "top-center": (0.5, 0),
        "top-right": (1, 0),
        "center": (0.5, 0.5),
        "bottom-left": (0, 1),
        "bottom-center": (0.5, 1),
        "bottom-right": (1, 1),
    }

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    def add_single(slide, px, py):
        if watermark_type == "text":
            txBox = slide.shapes.add_textbox(
                Emu(int(slide_w * px - slide_w * factor / 2)),
                Emu(int(slide_h * py - slide_h * factor / 2)),
                Emu(int(slide_w * factor)),
                Emu(int(slide_h * factor)),
            )
            _style_text_frame(txBox, text, font_size, opacity)
        else:
            _add_image_shape(slide, image_bytes, int(slide_w * px - slide_w * factor / 2), int(slide_h * py - slide_h * factor / 2), int(slide_w * factor))

    def add_tile(slide):
        step_x = int(slide_w * factor * 1.6)
        step_y = int(slide_h * factor * 1.6)
        if step_x <= 0 or step_y <= 0:
            return
        for ty in range(0, slide_h, step_y):
            for tx in range(0, slide_w, step_x):
                if watermark_type == "text":
                    box_w = min(int(slide_w * factor), slide_w - tx)
                    box_h = min(int(slide_h * factor), slide_h - ty)
                    txBox = slide.shapes.add_textbox(Emu(tx), Emu(ty), Emu(max(box_w, 1000)), Emu(max(box_h, 300)))
                    _style_text_frame(txBox, text, font_size, opacity)
                else:
                    _add_image_shape(slide, image_bytes, tx, ty, min(int(slide_w * factor), slide_w - tx))

    for idx in target_indices:
        slide = prs.slides[idx]
        if style == "tile":
            add_tile(slide)
        else:
            px, py = pos_map.get(position, (0.5, 0.5))
            add_single(slide, px, py)

    prs.save(output_path)


def _style_text_frame(txBox, text, font_size, opacity):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run._r.get_or_add_rPr()
        solidFill = run._r.rPr.makeelement(qn("a:solidFill"), {})
        srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": "000000"})
        solidFill.append(srgb)
        run._r.rPr.append(solidFill)
        alpha = solidFill.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))})
        solidFill.append(alpha)


def _add_image_shape(slide, image_bytes, left, top, img_w):
    import io
    from PIL import Image as PILImage
    img = PILImage.open(io.BytesIO(image_bytes))
    img_h = int(img_w * img.size[1] / img.size[0])
    slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, img_w, img_h)